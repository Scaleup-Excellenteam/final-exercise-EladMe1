import asyncio
import json
import os
import time
import openai
from pptx import Presentation
from dotenv import load_dotenv
from time import sleep

# Load environment variables from API.env
load_dotenv("API.env")

# Access the API key
API_KEY = os.getenv("API_KEY")

# 3 requests per minute
RATE_LIMIT = 3

# Variables to track rate limit
request_count = 0
last_request_time = 0


async def extract_text_from_powerpoint(filepath):
    """
    Extracts text from a PowerPoint presentation.

    Args:
        filepath (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of strings, where each string represents the text content of a slide.
              If an error occurs during extraction, an empty list is returned.
    """
    try:
        prs = Presentation(filepath)
        text_list = []
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text.append(paragraph.text)
            text_list.append(' '.join(text))
        return text_list
    except Exception as e:
        print(f"Error occurred while extracting text: {str(e)}")
        return []


async def ask_chatgpt(prompt):
    """
    Generates a response from the ChatGPT model based on a given prompt.

    Args:
        prompt (str): The prompt for generating the response.

    Returns:
        str: The generated response from the ChatGPT model.
    """
    global request_count, last_request_time

    # Check rate limit
    current_time = time.time()
    time_elapsed = current_time - last_request_time

    if time_elapsed < 60 and request_count >= RATE_LIMIT:
        sleep(60 - time_elapsed)
        request_count = 0
        last_request_time = current_time

    openai.api_key = API_KEY
    response = await asyncio.to_thread(openai.ChatCompletion.create,
                                       model="gpt-3.5-turbo",
                                       messages=[
                                           {"role": "user", "content": prompt}
                                       ],
                                       )

    request_count += 1
    last_request_time = time.time()

    if 'choices' in response:
        return response['choices'][0]['message']['content'].strip()
    return ""


def save_explanations_on_json_file(explanations, jason_file_name):
    """
        Saves the explanations to a JSON file.

        Args:
            explanations (list): List of explanations.
            filename (str): Name of the output JSON file.

        Return: None
        """
    json.dump(explanations, open(jason_file_name, "w"))


async def main():
    """
      Main function to extract text from PowerPoint, generate explanations using ChatGPT,
      and save the explanations to a JSON file.

      Args: None
      Return: None
      """

    #get the file pptx from user
    file_path = input("Enter the file path to the PowerPoint presentation: ")
    text = await extract_text_from_powerpoint(file_path)

    #loop on all the silde and send one slide each time to OpenAi Api
    #save the answer from OpenAi in list
    explanations = []
    for slide_text in text:
        prompt = f"Can you pls explain to me about: {slide_text}\n\n"
        response_text = await ask_chatgpt(prompt)
        explanations.append(response_text)
    print(explanations)

    #save the explanations in json file by name explanations
    explanationsFile = "explanations.json"
    save_explanations_on_json_file(explanations,explanationsFile)

if __name__ == '__main__':
    asyncio.run(main())
