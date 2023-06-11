from pptx import Presentation

async def extract_text_from_powerpoint(filepath):
    """
    Extracts text from a PowerPoint presentation.

    Args:
        filepath (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of strings, where each string represents the text content of a slide.
              If an error occurs during extraction, an empty list is returned.
    """
    try:
        prs = Presentation(filepath)
        text_list = []
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text.append(paragraph.text)
            text_list.append(' '.join(text))
        return text_list
    except Exception as e:
        print(f"Error occurred while extracting text: {str(e)}")
        return []
